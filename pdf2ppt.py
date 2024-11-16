import os
import time
import shutil
import argparse
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def ensure_directories():
    # Ensure "input" and "processed" directories exist
    os.makedirs("input", exist_ok=True)
    os.makedirs("processed", exist_ok=True)

def locate_pdf_file(pdf_path):
    # Check if the file exists in the root or input directory
    if os.path.exists(pdf_path):
        return pdf_path
    input_dir = os.path.join("input", os.path.basename(pdf_path))
    if os.path.exists(input_dir):
        return input_dir
    raise FileNotFoundError(f"File '{pdf_path}' not found in root or 'input' directory.")

def move_to_processed(pdf_path):
    # Move the PDF to the "processed" directory
    processed_dir = "processed"
    os.makedirs(processed_dir, exist_ok=True)
    shutil.move(pdf_path, os.path.join(processed_dir, os.path.basename(pdf_path)))

def pdf_to_ppt(pdf_path):
    # Generate a timestamped output directory
    timestamp = int(time.time())
    output_dir = f"{timestamp}_output_ppt"
    os.makedirs(output_dir, exist_ok=True)

    # Convert PDF pages to images
    images = convert_from_path(pdf_path)

    # Create a new PowerPoint presentation
    ppt = Presentation()

    # Add each image to a slide
    for image in images:
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank slide layout
        img_stream = image.convert("RGB")
        img_stream.save("temp.jpg", format="JPEG")
        slide.shapes.add_picture("temp.jpg", Inches(0), Inches(0), width=ppt.slide_width)

    # Save the presentation
    output_file = os.path.join(output_dir, "output.pptx")
    ppt.save(output_file)
    print(f"Presentation saved as {output_file}")

    # Move the PDF to the "processed" directory
    move_to_processed(pdf_path)

if __name__ == "__main__":
    # Set up argument parser
    parser = argparse.ArgumentParser(description="Convert a PDF to a PowerPoint presentation.")
    parser.add_argument("pdf_file", help="Path to the input PDF file")

    # Parse arguments
    args = parser.parse_args()

    # Ensure necessary directories and process the PDF
    try:
        ensure_directories()
        pdf_path = locate_pdf_file(args.pdf_file)
        pdf_to_ppt(pdf_path)
    except Exception as e:
        print(f"Error: {e}")
